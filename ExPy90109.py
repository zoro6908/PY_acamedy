# ExPy9xxxx 파이썬 개인 연습

# pandas를 이용한 주식 정보 처리

# lxml 인스톨 필요
# python-pptx 인스톨 필요

# 참고 자료
# https://ai-creator.tistory.com/307 : 메인
# https://wendys.tistory.com/173
# https://ai-creator.tistory.com/51
# https://wikidocs.net/981

# 한국 거래소 종목 코드 다운로드 : kind.krx.co.kr/corpgeneral/corpList.do?method=download

# Step1) 종목 코드 및 일별 시세 가져오기
import pandas as pd

def get_stock_code():
    # 한국 거래소 사이트에서 종목코드 다운로드
    stock_code = pd.read_html('http://kind.krx.co.kr/corpgeneral/corpList.do?method=download', header=0)[0]
    # 필요없는 컬럼 제외
    stock_code = stock_code[['회사명', '종목코드']]
    # 한글 컬럼명을 영문 컬럼명으로 변경
    stock_code = stock_code.rename(columns={'회사명':'company', '종목코드':'code'})
    # 종목코드가 6자리이기 때문에 6자지를 맞추어 줌
    stock_code.code = stock_code.code.map('{:06d}'.format)

    return stock_code


# df = pd.read_html('http://kind.krx.co.kr/corpgeneral/corpList.do?method=download', header=0)[0]
# print(df.head())

def get_daily_stock(code):
    df = pd.DataFrame()
    for page in range(1, 21):
        # 일별 시세
        url = 'http://finance.naver.com/item/sise_day.nhn?code={code}'.format(code=code)
        url = '{url}&page={page}'.format(url=url, page=page)
        # print(url)
        current_df = pd.read_html(url, header=0)[0]
        df = df.append(current_df, ignore_index=True)
    return df

# code = '005930' # 삼성전자 종목코드
# df = get_daily_stock(code)
# print(df)

def data_cleaning(df):
    # df.dropna()를 이용해 결측값 있는 행 제거
    df = df.dropna()

    # 한글로 된 컬럼명을 영어로 바꿔줌
    df = df.rename(columns={'날짜':'date', '종가':'close', '전일비':'diff', '시가':'open', '고가':'high',
                           '저가':'low', '거래량':'volume'})

    # 데이터의 타입을 int형으로 바꿔줌
    df[['close', 'diff', 'open', 'high', 'low', 'volume']] = \
        df[['close', 'diff', 'open', 'high', 'low', 'volume']].astype(int)

    # 컬럼명 'date'의 타입을 date로 바꿔줌
    df['date'] = pd.to_datetime(df['date'])

    # 일자(date)를 기준으로 오름차순 정렬
    df = df.sort_values(by=['date'], ascending=True)

    return df

#################################
## 함수 호출
#################################
# 종목 코드 가져오기
company = input('회사명을 입력하세요 : ')
stock_code = get_stock_code()

# 일별 시세 가져오기
code = stock_code[stock_code.company == company].code.values[0].strip()  ## strip() : 공백제거
df = get_daily_stock(code)

# 일별 시세 클린징
df = data_cleaning(df)

# print(df)

# Step2) 보고 자료 준비하기
import matplotlib.pyplot as plt
from pandas.plotting import table
import os

#################################
## 차트 그리기
#################################
plt.figure(figsize=(10, 4))
plt.plot(df['date'], df['close'])
plt.xlabel('date')
plt.ylabel('close')

#################################
## 차트 저장 및 출력하기
#################################
chart_fname = os.path.join("/Temp/stock_report", '{company}_chart.png'.format(company=company))
plt.savefig(chart_fname)
plt.show()

#################################
## 일별 시세 그리기
#################################
plt.figure(figsize=(15, 4))
ax = plt.subplot(111, frame_on=False)  # no visible frame
ax.xaxis.set_visible(False)  # hide the x axis
ax.yaxis.set_visible(False)  # hide the y axis
df = df.sort_values(by=['date'], ascending=False)
table(ax, df.head(10), loc='center', cellLoc='center', rowLoc='center')  # where df is your data frame

#################################
## 일별 시세 저장하기
#################################
table_fname = os.path.join("/Temp/stock_report", '{company}_table.png'.format(company=company))
plt.savefig(table_fname)

# Step3) 보고서 작성하기
import datetime
from pptx import Presentation  # 라이브러리
from pptx.util import Inches  # 사진, 표등을 그리기 위해
import os

#################################
## 파워포인트 객체 선언
#################################
today = datetime.datetime.today().strftime('%Y%m%d')
prs = Presentation()  # 파워포인트 객체 선언

#################################
## 제목 장표 추가
#################################
title_slide_layout = prs.slide_layouts[0]  # 0 : 제목슬라이드에 해당
slide = prs.slides.add_slide(title_slide_layout)  # 제목 슬라이드를 파워포인트 객체에 추가

# 제목 - 제목에 값넣기
title = slide.shapes.title  # 제목
title.text = "주식 보고서"  # 제목에 값 넣기
# 부제목
subtitle = slide.placeholders[1]  # 제목상자는 placeholders[0], 부제목상자는 [1]
subtitle.text = "보고서 작성일 : {date}".format(date=today)

#################################
## 차트 및 테이블 장표 추가
#################################
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)

shapes = slide.shapes
shapes.title.text = '{company}, {close}원에 거래 마감'.format(company=company, close=df.iloc[0]['close'])
print(shapes.title.text)

# 차트 추가
left = Inches(0.5)
height = Inches(2.5)
width = Inches(9)
top = Inches(2)
# width, hegith가 없을 경우 원본 사이즈로
pic = slide.shapes.add_picture(chart_fname, left, top, width=width, height=height)

# 테이블 추가
left = Inches(-1)
height = Inches(3)
width = Inches(12)
top = Inches(4)

pic = slide.shapes.add_picture(table_fname, left, top, width=width, height=height)
cursor_sp = slide.shapes[0]._element
cursor_sp.addprevious(pic._element)  # 해당 요소를 뒤로 보내기 합니다.

#################################
## 보고서 저장
#################################
ppt_fname = os.path.join("/Temp/stock_report", 'stock_report.pptx')
prs.save(ppt_fname)

# Step4) 보고서를 이메일로 전송하기
import smtplib
# 이메일 메시지에 다양한 형식을 중첩하여 담기 위한 객체
from email.mime.multipart import MIMEMultipart

# 이메일 메시지를 이진 데이터로 바꿔주는 인코더
from email import encoders

# 텍스트 형식
from email.mime.text import MIMEText
# 이미지 형식
from email.mime.image import MIMEImage
# 오디오 형식
from email.mime.audio import MIMEAudio

# 위의 모든 객체들을 생성할 수 있는 기본 객체
# MIMEBase(_maintype, _subtype)
# MIMEBase(<메인 타입>, <서브 타입>)
from email.mime.base import MIMEBase


#################################
## 함수 정의
#################################
def send_email(smtp_info, msg):
    with smtplib.SMTP(smtp_info["smtp_server"], smtp_info["smtp_port"]) as server:
        # TLS 보안 연결
        server.starttls()
        # 로그인
        server.login(smtp_info["smtp_user_id"], smtp_info["smtp_user_pw"])

        # 로그인 된 서버에 이메일 전송
        response = server.sendmail(msg['from'], msg['to'],
                                   msg.as_string())  # 메시지를 보낼때는 .as_string() 메소드를 사용해서 문자열로 바꿔줍니다.

        # 이메일을 성공적으로 보내면 결과는 {}
        if not response:
            print('이메일을 성공적으로 보냈습니다.')
        else:
            print(response)


def make_multimsg(msg_dict):
    multi = MIMEMultipart(_subtype='mixed')

    for key, value in msg_dict.items():
        # 각 타입에 적절한 MIMExxx()함수를 호출하여 msg 객체를 생성한다.
        if key == 'text':
            with open(value['filename'], encoding='utf-8') as fp:
                msg = MIMEText(fp.read(), _subtype=value['subtype'])
        elif key == 'image':
            with open(value['filename'], 'rb') as fp:
                msg = MIMEImage(fp.read(), _subtype=value['subtype'])
        elif key == 'audio':
            with open(value['filename'], 'rb') as fp:
                msg = MIMEAudio(fp.read(), _subtype=value['subtype'])
        else:
            with open(value['filename'], 'rb') as fp:
                msg = MIMEBase(value['maintype'], _subtype=value['subtype'])
                msg.set_payload(fp.read())
                encoders.encode_base64(msg)

        # 경로가 있는 경우, 파일의 이름만 추출 ex) res/stock_report/stock_report.pptx -> stock_report.pptx
        _, fname = os.path.split(value['filename'])
        # 파일 이름을 첨부파일 제목으로 추가
        msg.add_header('Content-Disposition', 'attachment', filename=fname)

        # 첨부파일 추가
        multi.attach(msg)

    return multi


#################################
## 함수 호출
#################################
smtp_info = dict({"smtp_server": "smtp.naver.com",  # SMTP 서버 주소
                  "smtp_user_id": "<송신자(sender) 메일 계정>@naver.com",
                  "smtp_user_pw": "<송신자(sender) 메일 패스워드>",
                  "smtp_port": 587})  # SMTP 서버 포트

msg_dict = {
    'application': {'maintype': 'application', 'subtype': 'octect-stream', 'filename': 'res/email_sending/test.pdf'}
    # 그외 첨부파일
}

# 메일 내용 작성
title = '({date}). 주식 보고서 분석 자료 입니다'.format(date=today)
content = '주식 보고서 분석 자료 입니다'
sender = "<송신자(sender) 메일 계정>@naver.com"
receiver = "<수신자(receiver) 메일 주소>@naver.com"
msg = MIMEText(_text=content, _charset="utf-8")

# 첨부파일 추가
msg_dict['application']['filename'] = ppt_fname
multi = make_multimsg(msg_dict)
multi['Subject'] = title
multi['From'] = sender
multi['To'] = receiver
multi.attach(msg)

# 이메일 전송
send_email(smtp_info, multi)
